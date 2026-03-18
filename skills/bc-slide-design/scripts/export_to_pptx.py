import os
import sys

# Bootstrap: ensure dependencies are installed (venv locally, pip in sandboxes)
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _bootstrap import ensure_deps
ensure_deps()

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
import copy
import re

# BC Brand colors as hex (oklch approximations)
THEMES = {
    'dark': {
        'bg': RGBColor(0x26, 0x23, 0x20),
        'text': RGBColor(0xED, 0xE8, 0xE0),
        'accent': RGBColor(0xC0, 0x30, 0x20),
        'accent2': RGBColor(0xD4, 0xC0, 0x8A),
        'muted': RGBColor(0x88, 0x82, 0x78),
        'font_display': 'Georgia',
        'font_body': 'Calibri',
    },
    'light': {
        'bg': RGBColor(0xF7, 0xF3, 0xED),
        'text': RGBColor(0x3A, 0x35, 0x30),
        'accent': RGBColor(0x8A, 0x10, 0x0B),
        'accent2': RGBColor(0xB2, 0x9D, 0x6C),
        'muted': RGBColor(0x78, 0x72, 0x68),
        'font_display': 'Georgia',
        'font_body': 'Calibri',
    }
}

# ─── Slide dimensions (16:9) ───────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ─── Layout constants ──────────────────────────────────────────────────────
MARGIN_LEFT = Inches(1)
MARGIN_RIGHT = Inches(1)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT  # ~11.333 in


# ─── Whitespace normalisation ──────────────────────────────────────────────

def _collapse_ws(text):
    """Collapse all whitespace (newlines, tabs, multiple spaces) into single
    spaces, matching how a browser renders HTML inline text."""
    if not text:
        return ''
    return ' '.join(text.split())


# ─── Rich-text helpers ─────────────────────────────────────────────────────

def _walk_inline(node, inherited_bold=False, inherited_italic=False,
                 inherited_color=None):
    """Yield (text, bold, italic, color) tuples by walking an element's
    children, preserving <strong>/<em>/<span> formatting.

    Each tuple represents one *run* worth of text.  Whitespace is collapsed
    at the run level; the caller must still strip leading/trailing space on
    the assembled line.
    """
    for child in node.children:
        if isinstance(child, NavigableString):
            text = _collapse_ws(str(child))
            if text:
                yield (text, inherited_bold, inherited_italic, inherited_color)
        elif isinstance(child, Tag):
            tag = child.name
            bold = inherited_bold
            italic = inherited_italic
            color = inherited_color

            if tag in ('strong', 'b'):
                bold = True
            if tag in ('em', 'i'):
                italic = True

            classes = child.get('class', [])
            if 'em-accent' in classes or 'em-gold' in classes:
                bold = True
                color = 'accent2'  # resolved later from theme

            # Recurse
            yield from _walk_inline(child, bold, italic, color)


def _add_runs_to_paragraph(para, element, theme, font_size=None,
                           font_name=None, default_color=None,
                           bold_override=None, italic_override=None):
    """Populate *para* with runs extracted from *element*.

    Handles <strong>, <em>, <span class="em-accent">, etc.
    Returns the paragraph for chaining.
    """
    if font_name is None:
        font_name = theme['font_body']
    if default_color is None:
        default_color = theme['text']

    runs = list(_walk_inline(element))

    # Collapse whitespace between adjacent runs (avoid doubled spaces at
    # run boundaries).
    texts = []
    for text, bold, italic, color in runs:
        texts.append(text)
    full_text = _collapse_ws(' '.join(texts))
    if not full_text:
        return para

    # Re-walk to create properly formatted runs.
    for idx, (text, bold, italic, color) in enumerate(runs):
        text = _collapse_ws(text)
        if not text:
            continue
        run = para.add_run()
        run.text = text
        run.font.size = font_size
        run.font.name = font_name
        if bold_override is not None:
            run.font.bold = bold_override
        else:
            run.font.bold = bold or None
        if italic_override is not None:
            run.font.italic = italic_override
        else:
            run.font.italic = italic or None

        if color == 'accent2':
            run.font.color.rgb = theme['accent2']
        else:
            run.font.color.rgb = default_color

    return para


def _add_plain_text_paragraph(tf, text, theme, font_size=Pt(18),
                              font_name=None, color=None,
                              alignment=None, bold=None, italic=None,
                              space_after=None, space_before=None,
                              is_first=False):
    """Add a simple plain-text paragraph (no rich-text element needed)."""
    if font_name is None:
        font_name = theme['font_body']
    if color is None:
        color = theme['text']

    if is_first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    run = p.add_run()
    run.text = _collapse_ws(text)
    run.font.size = font_size
    run.font.name = font_name
    run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if alignment is not None:
        p.alignment = alignment
    if space_after is not None:
        p.space_after = space_after
    if space_before is not None:
        p.space_before = space_before
    return p


# ─── Bullet XML helper ────────────────────────────────────────────────────

def _set_bullet(paragraph, bullet_char='\u2022', color=None):
    """Add a bullet marker to a python-pptx paragraph via oxml."""
    pPr = paragraph._p.get_or_add_pPr()
    buNone = pPr.find(qn('a:buNone'))
    if buNone is not None:
        pPr.remove(buNone)

    # Bullet character
    buChar = pPr.find(qn('a:buChar'))
    if buChar is None:
        buChar = pPr.makeelement(qn('a:buChar'), {})
        pPr.append(buChar)
    buChar.set('char', bullet_char)

    if color is not None:
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {
            'val': '%02X%02X%02X' % (color[0], color[1], color[2])
        })
        buClr.append(srgbClr)
        # Insert buClr before buChar so OOXML ordering is correct
        pPr.insert(list(pPr).index(buChar), buClr)

    # Bullet size (percentage of font size)
    buSzPct = pPr.find(qn('a:buSzPct'))
    if buSzPct is None:
        buSzPct = pPr.makeelement(qn('a:buSzPct'), {})
        pPr.insert(list(pPr).index(buChar), buSzPct)
    buSzPct.set('val', '100000')  # 100%


def _set_numbered_bullet(paragraph, start_at=1):
    """Add auto-numbered bullet to paragraph."""
    pPr = paragraph._p.get_or_add_pPr()
    buNone = pPr.find(qn('a:buNone'))
    if buNone is not None:
        pPr.remove(buNone)
    buAutoNum = pPr.makeelement(qn('a:buAutoNum'), {
        'type': 'arabicPeriod',
        'startAt': str(start_at),
    })
    pPr.append(buAutoNum)


def _set_indent(paragraph, level=0, hanging=Inches(0.3)):
    """Set paragraph indent + hanging indent for bullet alignment."""
    pPr = paragraph._p.get_or_add_pPr()
    margin = int(Inches(0.4) + level * Inches(0.4))
    pPr.set('marL', str(margin))
    pPr.set('indent', str(-int(hanging)))


# ─── Image helper ──────────────────────────────────────────────────────────

def _resolve_image_path(src, html_dir):
    """Resolve an <img> src to an absolute path, return None if not found."""
    if not src or src.startswith('data:') or src.startswith('http'):
        return None  # skip data URIs and remote URLs for now
    path = os.path.normpath(os.path.join(html_dir, src))
    if os.path.isfile(path):
        return path
    return None


def _add_image(slide, img_path, left, top, max_width, max_height):
    """Add an image to the slide, scaling to fit within max dimensions while
    maintaining aspect ratio."""
    try:
        pic = slide.shapes.add_picture(img_path, left, top)
    except Exception:
        return None  # skip if image can't be loaded

    # Scale to fit
    img_w = pic.width
    img_h = pic.height
    if img_w > max_width or img_h > max_height:
        ratio_w = max_width / img_w
        ratio_h = max_height / img_h
        ratio = min(ratio_w, ratio_h)
        pic.width = int(img_w * ratio)
        pic.height = int(img_h * ratio)

    return pic


# ─── Theme detection ───────────────────────────────────────────────────────

def detect_theme(soup):
    """Detect theme from data-theme attribute on <html>."""
    html_tag = soup.find('html')
    if html_tag and html_tag.get('data-theme') == 'light':
        return 'light'
    return 'dark'


# ─── Content element processing ────────────────────────────────────────────

def _process_elements(container, html_dir):
    """Walk *container* (a <section> or sub-div) and yield content tuples:

    ('heading', level, element)
    ('section-label', element)
    ('subtitle', element)
    ('para', element)
    ('bullet', element)              # <li> inside <ul>
    ('numbered', element, index)     # <li> inside <ol>
    ('prompt-block', element)
    ('callout', element)
    ('two-col', element)
    ('compare-row', element)
    ('image', resolved_path)
    """
    for child in container.children:
        if isinstance(child, NavigableString):
            continue
        if not isinstance(child, Tag):
            continue

        tag = child.name
        classes = child.get('class', [])

        # Skip navigation / non-content wrappers
        if tag in ('nav', 'script', 'style'):
            continue

        # Headings
        if tag == 'h1':
            yield ('heading', 1, child)
            continue
        if tag == 'h2':
            yield ('heading', 2, child)
            continue
        if tag == 'h3':
            yield ('heading', 3, child)
            continue

        # Section label
        if tag == 'p' and 'section-label' in classes:
            yield ('section-label', child)
            continue

        # Subtitle
        if tag == 'p' and 'subtitle' in classes:
            yield ('subtitle', child)
            continue

        # Prompt block
        if tag == 'div' and 'prompt-block' in classes:
            yield ('prompt-block', child)
            continue

        # Callout
        if tag == 'div' and 'callout' in classes:
            yield ('callout', child)
            continue

        # Two-column layout
        if tag == 'div' and 'two-col' in classes:
            yield ('two-col', child)
            continue

        # Compare row
        if tag == 'div' and 'compare-row' in classes:
            yield ('compare-row', child)
            continue

        # Unordered list
        if tag == 'ul':
            for li in child.find_all('li', recursive=False):
                yield ('bullet', li)
            continue

        # Ordered list
        if tag == 'ol':
            for idx, li in enumerate(child.find_all('li', recursive=False), 1):
                yield ('numbered', li, idx)
            continue

        # Image
        if tag == 'img':
            src = child.get('src', '')
            path = _resolve_image_path(src, html_dir)
            if path:
                yield ('image', path)
            continue

        # Plain paragraph
        if tag == 'p':
            yield ('para', child)
            continue

        # Dive into generic wrappers (div without special class, etc.)
        if tag == 'div':
            # Check for images inside
            for img in child.find_all('img', recursive=True):
                src = img.get('src', '')
                path = _resolve_image_path(src, html_dir)
                if path:
                    yield ('image', path)
            # Also yield any other recognized children
            yield from _process_elements(child, html_dir)


# ─── Slide builders ────────────────────────────────────────────────────────

def _set_slide_bg(slide, theme):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = theme['bg']


def _build_title_slide(prs, section, theme, html_dir):
    """Title slide: centered title + subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme)

    elements = list(_process_elements(section, html_dir))

    title_text = ''
    title_elem = None
    subtitle_text = ''
    subtitle_elem = None
    section_label_elem = None

    for item in elements:
        if item[0] == 'section-label' and section_label_elem is None:
            section_label_elem = item[1]
        elif item[0] == 'heading' and title_elem is None:
            title_elem = item[2]
        elif item[0] == 'subtitle' and subtitle_elem is None:
            subtitle_elem = item[1]
        elif item[0] == 'para' and subtitle_elem is None and title_elem is not None:
            subtitle_elem = item[1]

    cur_top = Inches(2.0)
    center_left = Inches(1.5)
    center_width = Inches(10.333)

    # Section label
    if section_label_elem is not None:
        txBox = slide.shapes.add_textbox(center_left, cur_top, center_width,
                                         Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_runs_to_paragraph(p, section_label_elem, theme,
                               font_size=Pt(14),
                               font_name=theme['font_body'],
                               default_color=theme['accent2'],
                               bold_override=True)
        cur_top += Inches(0.5)

    # Title
    if title_elem is not None:
        txBox = slide.shapes.add_textbox(center_left, cur_top, center_width,
                                         Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_runs_to_paragraph(p, title_elem, theme,
                               font_size=Pt(44),
                               font_name=theme['font_display'],
                               default_color=theme['text'])
        cur_top += Inches(2.0)

    # Subtitle
    if subtitle_elem is not None:
        txBox = slide.shapes.add_textbox(center_left, cur_top, center_width,
                                         Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_runs_to_paragraph(p, subtitle_elem, theme,
                               font_size=Pt(24),
                               font_name=theme['font_body'],
                               default_color=theme['muted'])

    return slide


def _build_split_slide(prs, section, theme, html_dir):
    """Split slide: two-column layout (text + image)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme)

    # Determine column order: which side has the image?
    children = [c for c in section.children if isinstance(c, Tag)]
    image_left = False
    image_path = None
    text_div = None

    for child in children:
        classes = child.get('class', [])
        if 'slide-image' in classes:
            img = child.find('img')
            if img:
                src = img.get('src', '')
                image_path = _resolve_image_path(src, html_dir)
            if not text_div:
                image_left = True
        elif 'slide-text' in classes:
            text_div = child

    if text_div is None:
        # Fallback: treat whole section as text
        text_div = section

    # Layout constants for split
    half_width = Inches(5.5)
    col_margin = Inches(0.8)

    if image_left:
        img_left = col_margin
        text_left = Inches(6.8)
    else:
        text_left = col_margin
        img_left = Inches(6.8)

    # Image side
    if image_path:
        _add_image(slide, image_path, img_left, Inches(0.8),
                   half_width, Inches(5.8))

    # Text side
    _render_content_to_area(slide, text_div, theme, html_dir,
                            left=text_left, top=Inches(0.6),
                            width=half_width, max_height=Inches(6.2))

    return slide


def _build_centered_slide(prs, section, theme, html_dir):
    """Centered slide: content centered on page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme)

    _render_content_to_area(slide, section, theme, html_dir,
                            left=Inches(1.5), top=Inches(0.6),
                            width=Inches(10.333), max_height=Inches(6.4),
                            center=True)
    return slide


def _build_content_slide(prs, section, theme, html_dir):
    """Default content slide: left-aligned title + body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme)

    _render_content_to_area(slide, section, theme, html_dir,
                            left=MARGIN_LEFT, top=Inches(0.4),
                            width=CONTENT_WIDTH, max_height=Inches(6.6))
    return slide


# ─── Content renderer ──────────────────────────────────────────────────────

def _render_content_to_area(slide, container, theme, html_dir,
                            left, top, width, max_height,
                            center=False):
    """Render all recognised elements from *container* into a text box at
    the given position.  Images are placed separately as picture shapes.

    This is the workhorse that handles headings, paragraphs, lists,
    callouts, prompt blocks, two-col layouts, etc.
    """
    elements = list(_process_elements(container, html_dir))

    # We'll build a single text frame for text content and place images
    # alongside.
    txBox = slide.shapes.add_textbox(left, top, width, max_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    first_para = True
    cur_image_top = top  # track where to put images
    heading_seen = False
    ol_counter = 1

    for item in elements:
        kind = item[0]

        # Get or create paragraph
        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()

        # ── Section label ──────────────────────────────────────────
        if kind == 'section-label':
            elem = item[1]
            if center:
                p.alignment = PP_ALIGN.CENTER
            _add_runs_to_paragraph(p, elem, theme,
                                   font_size=Pt(12),
                                   font_name=theme['font_body'],
                                   default_color=theme['accent2'],
                                   bold_override=True)
            p.space_after = Pt(2)
            p.space_before = Pt(6)
            continue

        # ── Headings ───────────────────────────────────────────────
        if kind == 'heading':
            level = item[1]
            elem = item[2]
            sizes = {1: Pt(36), 2: Pt(30), 3: Pt(22)}
            font_size = sizes.get(level, Pt(22))
            if center:
                p.alignment = PP_ALIGN.CENTER
            _add_runs_to_paragraph(p, elem, theme,
                                   font_size=font_size,
                                   font_name=theme['font_display'],
                                   default_color=theme['accent'])
            p.space_after = Pt(4)
            if heading_seen:
                p.space_before = Pt(16)

            # Add accent line after first heading (h1 or h2)
            if not heading_seen and level <= 2 and not center:
                heading_seen = True
                line_p = tf.add_paragraph()
                line_p.space_after = Pt(8)
                # We'll add a thin accent-colored run as a visual separator
                run = line_p.add_run()
                run.text = '\u2500' * 16  # box-drawing horizontal line
                run.font.size = Pt(10)
                run.font.color.rgb = theme['accent2']

            heading_seen = True
            continue

        # ── Subtitle ───────────────────────────────────────────────
        if kind == 'subtitle':
            elem = item[1]
            if center:
                p.alignment = PP_ALIGN.CENTER
            _add_runs_to_paragraph(p, elem, theme,
                                   font_size=Pt(20),
                                   font_name=theme['font_body'],
                                   default_color=theme['muted'])
            p.space_after = Pt(12)
            continue

        # ── Paragraph ──────────────────────────────────────────────
        if kind == 'para':
            elem = item[1]
            if center:
                p.alignment = PP_ALIGN.CENTER
            _add_runs_to_paragraph(p, elem, theme,
                                   font_size=Pt(18),
                                   font_name=theme['font_body'],
                                   default_color=theme['text'])
            p.space_after = Pt(10)
            continue

        # ── Bullet (unordered list item) ───────────────────────────
        if kind == 'bullet':
            elem = item[1]
            _add_runs_to_paragraph(p, elem, theme,
                                   font_size=Pt(18),
                                   font_name=theme['font_body'],
                                   default_color=theme['text'])
            _set_bullet(p, '\u2022',
                        color=(theme['accent2'][0],
                               theme['accent2'][1],
                               theme['accent2'][2]))  # accent-colored bullets
            _set_indent(p, level=0)
            p.space_after = Pt(6)
            continue

        # ── Numbered list item ─────────────────────────────────────
        if kind == 'numbered':
            elem = item[1]
            idx = item[2]
            _add_runs_to_paragraph(p, elem, theme,
                                   font_size=Pt(18),
                                   font_name=theme['font_body'],
                                   default_color=theme['text'])
            _set_numbered_bullet(p, start_at=idx)
            _set_indent(p, level=0)
            p.space_after = Pt(6)
            continue

        # ── Prompt block ───────────────────────────────────────────
        if kind == 'prompt-block':
            elem = item[1]
            # Render as indented italic block
            p.space_before = Pt(8)
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', str(int(Inches(0.6))))

            # Extract all text from the prompt block
            _add_runs_to_paragraph(p, elem, theme,
                                   font_size=Pt(16),
                                   font_name=theme['font_body'],
                                   default_color=theme['muted'],
                                   italic_override=True)
            p.space_after = Pt(8)

            # Add a left-border indicator character at the start
            # (shift existing runs right, add a colored bar run first)
            if len(p.runs) > 0:
                bar_run = p.runs[0]
                # Prepend a visual indicator
                bar_run.text = '\u275D ' + bar_run.text
            continue

        # ── Callout ────────────────────────────────────────────────
        if kind == 'callout':
            elem = item[1]
            p.space_before = Pt(12)
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', str(int(Inches(0.5))))
            if center:
                p.alignment = PP_ALIGN.CENTER

            _add_runs_to_paragraph(p, elem, theme,
                                   font_size=Pt(20),
                                   font_name=theme['font_display'],
                                   default_color=theme['accent2'])
            p.space_after = Pt(12)
            continue

        # ── Two-column layout ──────────────────────────────────────
        if kind == 'two-col':
            elem = item[1]
            cols = elem.find_all('div', class_='col', recursive=False)
            if not cols:
                cols = elem.find_all('div', recursive=False)

            for col_idx, col in enumerate(cols):
                # Column header
                header = col.find(class_='col-header')
                if header is None:
                    header = col.find(['h3', 'h4'])
                if header is not None:
                    if col_idx > 0 or not first_para:
                        p = tf.add_paragraph()
                    p.space_before = Pt(10)
                    _add_runs_to_paragraph(p, header, theme,
                                           font_size=Pt(18),
                                           font_name=theme['font_display'],
                                           default_color=theme['accent'],
                                           bold_override=True)
                    p.space_after = Pt(4)

                # Column list items
                for li in col.find_all('li'):
                    p = tf.add_paragraph()
                    _add_runs_to_paragraph(p, li, theme,
                                           font_size=Pt(16),
                                           font_name=theme['font_body'],
                                           default_color=theme['text'])
                    _set_bullet(p, '\u2022',
                                color=(theme['accent2'][0],
                                       theme['accent2'][1],
                                       theme['accent2'][2]))
                    _set_indent(p, level=0)
                    p.space_after = Pt(4)

                # Column paragraphs (if no list items)
                if not col.find('li'):
                    for cp in col.find_all('p'):
                        p = tf.add_paragraph()
                        _add_runs_to_paragraph(p, cp, theme,
                                               font_size=Pt(16),
                                               font_name=theme['font_body'],
                                               default_color=theme['text'])
                        p.space_after = Pt(4)

            first_para = False
            continue

        # ── Compare row ────────────────────────────────────────────
        if kind == 'compare-row':
            elem = item[1]
            items = elem.find_all('div', class_='compare-item')
            if not items:
                items = elem.find_all('div', recursive=False)

            for ci in items:
                label_span = ci.find('span', class_='compare-label')
                text_span = ci.find('span', class_='compare-text')

                p = tf.add_paragraph()
                first_para = False

                if label_span:
                    run = p.add_run()
                    run.text = _collapse_ws(label_span.get_text()) + '  '
                    run.font.size = Pt(16)
                    run.font.name = theme['font_body']
                    run.font.bold = True
                    run.font.color.rgb = theme['accent2']

                if text_span:
                    run = p.add_run()
                    run.text = _collapse_ws(text_span.get_text())
                    run.font.size = Pt(16)
                    run.font.name = theme['font_body']
                    run.font.color.rgb = theme['text']

                # If no spans, just dump all text
                if not label_span and not text_span:
                    _add_runs_to_paragraph(p, ci, theme,
                                           font_size=Pt(16),
                                           font_name=theme['font_body'],
                                           default_color=theme['text'])

                p.space_after = Pt(6)
            continue

        # ── Image ──────────────────────────────────────────────────
        if kind == 'image':
            img_path = item[1]
            # Place image below current text area
            img_top = top + Inches(0.5)  # approximate; placed independently
            _add_image(slide, img_path, left, img_top,
                       width, max_height - Inches(1))
            # Add an empty paragraph as spacer in the text frame
            first_para = False
            continue

    return txBox


# ─── Main extraction + build pipeline ──────────────────────────────────────

def extract_and_build(html_path, output_path):
    """Parse the HTML and build PPTX directly (no intermediate dict)."""
    html_dir = os.path.dirname(os.path.abspath(html_path))

    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')

    theme_name = detect_theme(soup)
    theme = THEMES[theme_name]

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for section in soup.find_all('section', class_='slide'):
        classes = section.get('class', [])

        if 'title-slide' in classes:
            _build_title_slide(prs, section, theme, html_dir)
        elif 'split-slide' in classes:
            _build_split_slide(prs, section, theme, html_dir)
        elif 'centered-slide' in classes:
            _build_centered_slide(prs, section, theme, html_dir)
        else:
            _build_content_slide(prs, section, theme, html_dir)

    prs.save(output_path)
    return output_path


def export(html_path, output_path=None):
    """Main export function."""
    if output_path is None:
        output_path = html_path.rsplit('.', 1)[0] + '.pptx'

    extract_and_build(html_path, output_path)
    return output_path


if __name__ == '__main__':
    import sys
    if len(sys.argv) < 2:
        print("Usage: python export_to_pptx.py <presentation.html> [output.pptx]")
        sys.exit(1)

    html_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None
    result = export(html_path, output_path)
    print(f"Exported to: {result}")
